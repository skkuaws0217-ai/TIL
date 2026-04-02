#!/usr/bin/env python3
"""프로젝트 PPT 자료 생성.

모델 구조도, 데이터 흐름도, AWS 아키텍처, 사용법을 포함한
프레젠테이션 자료를 생성한다.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # 배경색
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0B, 0x1D, 0x3A)

    # 제목
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # 부제
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x8E, 0xB8, 0xE5)
    p2.alignment = PP_ALIGN.CENTER
    return slide


def add_content_slide(prs, title, content_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # 제목 바
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0B, 0x1D, 0x3A)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)

    # 본문
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # 들여쓰기 처리
        indent = 0
        text = line
        if line.startswith("    "):
            indent = 1
            text = line.strip()
        if line.startswith("        "):
            indent = 2
            text = line.strip()

        p.text = text
        p.level = indent

        if line.startswith("##"):
            p.text = line.replace("## ", "").replace("##", "")
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x0B, 0x1D, 0x3A)
        elif line.startswith("["):
            p.font.size = Pt(11)
            p.font.italic = True
            p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        else:
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(4)

    return slide


def add_diagram_slide(prs, title, boxes, arrows_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 제목
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0B, 0x1D, 0x3A)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tf.margin_left = Inches(0.5)

    # 박스들 배치
    colors = [
        RGBColor(0x1B, 0x5E, 0x20),  # green
        RGBColor(0x0D, 0x47, 0xA1),  # blue
        RGBColor(0xE6, 0x51, 0x00),  # orange
        RGBColor(0x4A, 0x14, 0x8C),  # purple
        RGBColor(0xB7, 0x1C, 0x1C),  # red
        RGBColor(0x00, 0x69, 0x5C),  # teal
    ]

    y_start = 1.3
    for i, (box_title, box_content) in enumerate(boxes):
        color = colors[i % len(colors)]
        left = Inches(0.3 + (i % 3) * 3.2)
        top = Inches(y_start + (i // 3) * 2.4)
        width = Inches(3.0)
        height = Inches(2.0)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = box_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        for line in box_content.split("\n"):
            p2 = tf.add_paragraph()
            p2.text = line
            p2.font.size = Pt(10)
            p2.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
            p2.space_before = Pt(2)

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: 표지 ─────────────────────────────────────────
    add_title_slide(
        prs,
        "AI 기반 폐질환 및 희귀질환\n진단 보조 프로그램",
        "4-Phase Diagnostic Pipeline  |  536 Diseases  |  126 Reference Items\nSKKU Pre-Project  |  2026-04"
    )

    # ── Slide 2: 프로젝트 개요 ────────────────────────────────
    add_content_slide(prs, "프로젝트 개요", [
        "## 목적",
        "MIMIC-IV 데이터와 의학 참조 데이터 기반 AI 폐질환 진단 보조 도구",
        "",
        "## 4단계 파이프라인",
        "    Phase 1: X-ray 이미지 AI 분석 (CheXNet DenseNet121)",
        "    Phase 2: Lab + Vitals/Respiratory/Hemodynamic + 미생물 + 증상 다중모달 매칭",
        "    Phase 3: 376개 희귀 폐질환 HPO 빈도가중 스크리닝",
        "    Phase 4: 임상소견서 자동 생성 (Bedrock Claude / 로컬 템플릿)",
        "",
        "## 질환 커버리지",
        "    일반 폐질환 82개 + 기타 폐관련 70개 + 희귀 376개 + YAML 상세 17개 = 536개",
        "",
        "## 참조 데이터",
        "    Lab 89개 항목 (lab_reference_ranges_v3.yaml)",
        "    Vitals/Respiratory/Hemodynamic 37개 항목 (vitals_respiratory_hemodynamic_reference_range_v1.yaml)",
        "    7개 파일이 유일한 기준 출처",
    ])

    # ── Slide 3: 4단계 파이프라인 흐름도 ──────────────────────
    add_diagram_slide(prs, "4-Phase Diagnostic Pipeline", [
        ("Phase 1: X-ray AI", "CheXNet (DenseNet121)\n14 CheXpert labels\nGradCAM 히트맵\nAI 영상 키워드 51개\n→ 1차 질환 후보"),
        ("Phase 2: Multi-modal", "Lab (89항목) 해석\nVRH (37항목) 해석\n미생물 매칭\nHPO 증상 매칭\n가중 스코어링 (S/L/R/M)"),
        ("Phase 3: Rare Disease", "376개 희귀질환\nHPO 빈도가중 매칭\n3,468 HPO 레코드\n유전자 검사 제안\n확진검사 계획"),
        ("Phase 4: Report", "Bedrock Claude (AWS)\n또는 로컬 템플릿\n한국어 임상소견서\nMarkdown / JSON"),
        ("Knowledge Base", "7개 핵심 파일\nYAML 4개 + Excel 3개\n536개 질환 레지스트리\n1,416 HPO IDs"),
        ("입력 방법", "JSON 파일\nLab PDF 자동 파싱\n대화형 CLI\nFastAPI REST API"),
    ])

    # ── Slide 4: 데이터 흐름도 ────────────────────────────────
    add_content_slide(prs, "데이터 흐름도", [
        "## 환자 데이터 입력",
        "    JSON 직접 입력 / Lab PDF 자동 파싱 / 대화형 CLI / REST API",
        "",
        "## Phase 1 흐름 (X-ray 제공 시)",
        "    X-ray 이미지 → 전처리(224x224) → CheXNet → 14 labels 확률",
        "    → CheXpert→키워드 매핑 → 528개 질환 DB '영상 키워드' 매칭 → 1차 후보",
        "",
        "## Phase 2 흐름",
        "    Lab 값 → lab_reference_ranges_v3.yaml 기준 판정 → medical_term + disease_associations",
        "    VRH 값 → vitals_reference_v1.yaml 기준 판정 → thresholds + scoring(NEWS2/qSOFA/CURB-65/PESI)",
        "    미생물 → Excel DB 미생물 소견 + YAML micro_findings 매칭",
        "    증상 → HPO 코드 변환 → 질환별 HPO 리스트 매칭",
        "    → 질환별 S/L/R/M 가중치로 종합 스코어 → 순위 산출",
        "",
        "## Phase 3 흐름 (트리거 시)",
        "    환자 HPO → 376개 희귀질환 3,468 HPO 대비 빈도가중 매칭",
        "    → IC(Information Content) 보정 → 순위 → 유전자 검사 제안",
        "",
        "## Phase 4 흐름",
        "    Phase 1~3 결과 JSON → Bedrock Claude 또는 템플릿 → 임상소견서",
    ])

    # ── Slide 5: 스코어링 알고리즘 ────────────────────────────
    add_content_slide(prs, "스코어링 알고리즘 (Fact-Based)", [
        "## Phase 2: 가중 다중모달 스코어",
        "    score = Sum(w_c x match_ratio_c) / Sum(active_w_c)",
        "    w = 질환별 S(증상)/L(Lab)/R(영상)/M(미생물) 가중치 (합 = 1.0)",
        "",
        "## 가중치 출처",
        "    명시적 가중치 58개 질환: YAML/Excel 기재값 (Harrison's, GOLD, ATS/IDSA 등)",
        "    기본값 (일반): S:0.25 L:0.20 R:0.35 M:0.20 (YAML 17개 가중 평균)",
        "    기본값 (희귀): S:0.45 L:0.20 R:0.20 M:0.15 (HPO 중심 접근)",
        "",
        "## 보정 계수",
        "    Critical Lab 보너스: +0.05 [CAP Critical Values]",
        "    NEWS2 >= 7 + 감염성 질환: +0.03 [RCP NEWS2 2017]",
        "    다중모달 커버리지: sqrt(active/available) [Bayesian convergence]",
        "    최소 기준 수 정규화: min 3개 [과대평가 방지]",
        "",
        "## Phase 3: HPO 빈도가중 스코어",
        "    Obligate(1.0) > Very frequent(0.8) > Frequent(0.6) > Occasional(0.3) > Very rare(0.1)",
        "    IC 보정: -log2(diseases_with_hpo / total) x 0.01",
        "[Kohler et al. AJHG 2009; Orphanet phenotype-driven approach]",
    ])

    # ── Slide 6: AWS 아키텍처 ─────────────────────────────────
    add_diagram_slide(prs, "AWS 아키텍처", [
        ("Amazon SageMaker", "CheXNet 모델 배포\nml.g4dn.xlarge (GPU)\nReal-time Endpoint\n→ Phase 1 X-ray 추론"),
        ("Amazon Bedrock", "Claude Sonnet\n임상소견서 자연어 생성\n→ Phase 4 Report\n한국어 의학 용어"),
        ("Amazon S3", "참조 데이터 (YAML/Excel)\nX-ray 이미지 저장\nLab PDF 저장\n소견서 결과 저장"),
        ("API Gateway + Lambda", "REST API 엔드포인트\nPOST /api/v1/diagnose\nGET /api/v1/health\n또는 ECS Fargate"),
        ("Amazon CloudWatch", "로깅 + 모니터링\nSageMaker 추론 지연\nBedrock 호출 추적\n오류 알림"),
        ("환경변수 전환", "LUNG_DX_XRAY_BACKEND\nLUNG_DX_REPORT_BACKEND\nLUNG_DX_DATA_BACKEND\n코드 수정 없이 전환"),
    ])

    # ── Slide 7: AWS 마이그레이션 단계 ────────────────────────
    add_content_slide(prs, "AWS 마이그레이션 로드맵", [
        "## Stage 1: 로컬 개발 (현재)",
        "    CheXNet on Apple MPS / CPU",
        "    템플릿 기반 소견서",
        "    로컬 파일시스템",
        "",
        "## Stage 2: Bedrock 연동",
        "    aws configure → LUNG_DX_REPORT_BACKEND=bedrock",
        "    Claude Sonnet으로 자연어 소견서 생성",
        "",
        "## Stage 3: S3 데이터",
        "    YAML/Excel/이미지를 S3 업로드",
        "    LUNG_DX_DATA_BACKEND=s3",
        "",
        "## Stage 4: SageMaker 모델 배포",
        "    CheXNet → model.tar.gz → SageMaker Endpoint",
        "    LUNG_DX_XRAY_BACKEND=sagemaker",
        "",
        "## Stage 5: 풀 클라우드",
        "    ECS Fargate 또는 Lambda + API Gateway",
        "    CloudWatch 모니터링",
        "    CDK 인프라 코드",
    ])

    # ── Slide 8: 사용법 요약 ──────────────────────────────────
    add_content_slide(prs, "사용법 요약", [
        "## CLI - JSON 파일 입력",
        "    python -m lung_dx.cli --json patient.json",
        "",
        "## CLI - Lab PDF 자동 파싱",
        "    python -m lung_dx.cli --json patient.json --lab-pdf lab_results.pdf",
        "    (한국어/영문 검사명 자동 인식 → YAML ItemID 매칭)",
        "",
        "## CLI - 대화형 모드",
        "    python -m lung_dx.cli --interactive",
        "",
        "## FastAPI 웹 서버",
        "    uvicorn lung_dx.main:app --reload --port 8000",
        "    Swagger UI: http://localhost:8000/docs",
        "    POST /api/v1/diagnose → 진단 실행",
        "",
        "## 새 환자 검증",
        "    1) sample_patient.json 형식으로 JSON 작성",
        "    2) python -m lung_dx.cli --json my_patient.json 실행",
        "    3) 터미널에 임상소견서 출력 + .result.json 저장",
        "",
        "## 분석 기준",
        "    어떤 입력 경로든 동일한 YAML 기준으로 판정 (MIMIC-IV 종속 아님)",
    ])

    # ── Slide 9: 핵심 수치 ────────────────────────────────────
    add_content_slide(prs, "핵심 수치 요약", [
        "## 질환 데이터베이스",
        "    전체 536개 질환 (일반 82 + 기타 70 + 희귀 376 + YAML 17)",
        "    1,416개 고유 HPO ID",
        "    51개 AI 영상 키워드",
        "    유전자 정보 보유 101개 희귀질환",
        "",
        "## 검사 참조 항목",
        "    Lab: 89개 (MIMIC 53 + 외부 36), 13개 카테고리 (A~M)",
        "    VRH: 37개 (Vital Signs + Respiratory + Hemodynamic)",
        "    Critical 값: 20개 Lab 항목",
        "    Threshold: 15개 Lab + 다수 VRH 항목",
        "",
        "## 스코어링 시스템",
        "    NEWS2, qSOFA, CURB-65, PESI, Wells PE Score",
        "    파생 지표: S/F ratio, P/F ratio, Driving Pressure",
        "",
        "## 코드 규모",
        "    Python 파일 52개, 코드 6,187줄",
    ])

    # 저장
    output_path = "lung_dx_presentation.pptx"
    prs.save(output_path)
    print(f"PPT 생성 완료: {output_path}")


if __name__ == "__main__":
    main()
